from __future__ import annotations

import os
from datetime import datetime
from typing import Any

import pandas as pd

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# ==========================================================
# Theme
# ==========================================================

PRIMARY = RGBColor(37, 99, 235)
DARK = RGBColor(15, 23, 42)
LIGHT = RGBColor(248, 250, 252)
GRAY = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
SUCCESS = RGBColor(22, 163, 74)


# ==========================================================
# Presentation helpers
# ==========================================================

def create_presentation() -> Presentation:
    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    return prs


def add_background(slide):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = LIGHT


def add_header(slide, title: str):

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.33),
        Inches(0.55)
    )

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY

    shape.line.color.rgb = PRIMARY

    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE


def add_footer(slide, slide_number: int):

    textbox = slide.shapes.add_textbox(
        Inches(0.4),
        Inches(7.05),
        Inches(12.5),
        Inches(0.25)
    )

    tf = textbox.text_frame

    p = tf.paragraphs[0]

    p.text = (
        "AI Business Intelligence Platform"
        f"        Slide {slide_number}"
    )

    p.font.size = Pt(10)
    p.font.color.rgb = GRAY


def add_title(slide, text):

    textbox = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.75),
        Inches(12),
        Inches(0.5)
    )

    tf = textbox.text_frame

    p = tf.paragraphs[0]

    p.text = text

    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK


def add_subtitle(slide, text):

    textbox = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(1.25),
        Inches(12),
        Inches(0.35)
    )

    tf = textbox.text_frame

    p = tf.paragraphs[0]

    p.text = text

    p.font.size = Pt(13)
    p.font.color.rgb = GRAY


# ==========================================================
# Card
# ==========================================================

def add_metric_card(
    slide,
    left,
    top,
    width,
    height,
    title,
    value
):

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )

    card.fill.solid()
    card.fill.fore_color.rgb = WHITE

    card.line.color.rgb = PRIMARY

    tf = card.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY

    p2 = tf.add_paragraph()
    p2.text = str(value)
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = DARK


# ==========================================================
# Utility
# ==========================================================

def dataset_metrics(df: pd.DataFrame):

    rows = len(df)

    cols = len(df.columns)

    missing = int(
        df.isna()
        .sum()
        .sum()
    )

    duplicates = int(
        df.duplicated()
        .sum()
    )

    total_cells = max(
        1,
        rows * cols
    )

    quality = round(
        (
            1
            - missing / total_cells
        ) * 100,
        2
    )

    return {
        "rows": rows,
        "columns": cols,
        "missing": missing,
        "duplicates": duplicates,
        "quality": quality
    }


# ==========================================================
# Slide 1
# ==========================================================

def build_cover_slide(
    prs,
    filename
):

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    add_header(
        slide,
        "Executive Business Intelligence Report"
    )

    title = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(1)
    )

    tf = title.text_frame

    p = tf.paragraphs[0]

    p.text = "Executive BI Report"

    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = DARK

    subtitle = tf.add_paragraph()

    subtitle.text = (
        f"Dataset : {filename}"
    )

    subtitle.font.size = Pt(18)
    subtitle.font.color.rgb = PRIMARY

    date_box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(3.1),
        Inches(6),
        Inches(1)
    )

    dt = date_box.text_frame

    p = dt.paragraphs[0]

    p.text = (
        "Generated :\n"
        + datetime.now().strftime(
            "%d %B %Y\n%I:%M %p"
        )
    )

    p.font.size = Pt(18)
    p.font.color.rgb = GRAY

    footer = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(5.8),
        Inches(7),
        Inches(0.6)
    )

    ft = footer.text_frame

    p = ft.paragraphs[0]

    p.text = (
        "Generated by\n"
        "AI Business Intelligence Platform"
    )

    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS

    add_footer(
        slide,
        1
    )


# ==========================================================
# Slide 2
# ==========================================================

def build_overview_slide(
    prs,
    df
):

    metrics = dataset_metrics(df)

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    add_header(
        slide,
        "Dataset Overview"
    )

    add_title(
        slide,
        "Dataset Summary"
    )

    add_subtitle(
        slide,
        "High-level information about the uploaded dataset."
    )

    add_metric_card(
        slide,
        Inches(0.6),
        Inches(2.0),
        Inches(2.3),
        Inches(1.3),
        "Rows",
        metrics["rows"]
    )

    add_metric_card(
        slide,
        Inches(3.2),
        Inches(2.0),
        Inches(2.3),
        Inches(1.3),
        "Columns",
        metrics["columns"]
    )

    add_metric_card(
        slide,
        Inches(5.8),
        Inches(2.0),
        Inches(2.3),
        Inches(1.3),
        "Missing",
        metrics["missing"]
    )

    add_metric_card(
        slide,
        Inches(8.4),
        Inches(2.0),
        Inches(2.3),
        Inches(1.3),
        "Duplicates",
        metrics["duplicates"]
    )

    add_metric_card(
        slide,
        Inches(11.0),
        Inches(2.0),
        Inches(1.9),
        Inches(1.3),
        "Quality",
        f"{metrics['quality']}%"
    )

    add_footer(
        slide,
        2
    )
    # ==========================================================
# Table Helpers
# ==========================================================

def add_table_title(
    slide,
    text,
    top
):
    textbox = slide.shapes.add_textbox(
        Inches(0.6),
        top,
        Inches(12),
        Inches(0.4)
    )

    tf = textbox.text_frame

    p = tf.paragraphs[0]

    p.text = text

    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK


def add_simple_table(
    slide,
    data,
    left,
    top,
    width,
    row_height=0.38
):
    rows = len(data)
    cols = len(data[0])

    table = slide.shapes.add_table(
        rows,
        cols,
        left,
        top,
        width,
        Inches(row_height * rows)
    ).table

    column_width = width / cols

    for c in range(cols):
        table.columns[c].width = column_width

    for r in range(rows):

        table.rows[r].height = Inches(row_height)

        for c in range(cols):

            cell = table.cell(r, c)

            cell.text = str(data[r][c])

            paragraph = cell.text_frame.paragraphs[0]

            paragraph.font.size = Pt(11)

            paragraph.alignment = PP_ALIGN.CENTER

            if r == 0:

                paragraph.font.bold = True

                paragraph.font.color.rgb = WHITE

                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY

            else:

                paragraph.font.color.rgb = DARK

                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE


# ==========================================================
# Numeric Summary
# ==========================================================

def build_numeric_summary(
    df
):
    numeric = df.select_dtypes(
        include="number"
    )

    rows = [
        [
            "Column",
            "Mean",
            "Minimum",
            "Maximum",
            "Sum"
        ]
    ]

    for column in numeric.columns[:10]:

        series = numeric[column].dropna()

        if len(series) == 0:
            continue

        rows.append([
            column,
            f"{series.mean():,.2f}",
            f"{series.min():,.2f}",
            f"{series.max():,.2f}",
            f"{series.sum():,.2f}"
        ])

    return rows


# ==========================================================
# Slide 3
# Executive KPI Summary
# ==========================================================

def build_kpi_slide(
    prs,
    df
):

    metrics = dataset_metrics(df)

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    add_header(
        slide,
        "Executive KPI Dashboard"
    )

    add_title(
        slide,
        "Executive KPI Summary"
    )

    add_subtitle(
        slide,
        "Important business metrics generated from the uploaded dataset."
    )

    add_metric_card(
        slide,
        Inches(0.8),
        Inches(2),
        Inches(2.2),
        Inches(1.4),
        "Rows",
        metrics["rows"]
    )

    add_metric_card(
        slide,
        Inches(3.3),
        Inches(2),
        Inches(2.2),
        Inches(1.4),
        "Columns",
        metrics["columns"]
    )

    add_metric_card(
        slide,
        Inches(5.8),
        Inches(2),
        Inches(2.2),
        Inches(1.4),
        "Quality",
        f"{metrics['quality']}%"
    )

    add_metric_card(
        slide,
        Inches(8.3),
        Inches(2),
        Inches(2.2),
        Inches(1.4),
        "Missing",
        metrics["missing"]
    )

    add_metric_card(
        slide,
        Inches(10.8),
        Inches(2),
        Inches(2.0),
        Inches(1.4),
        "Duplicates",
        metrics["duplicates"]
    )

    numeric = len(
        df.select_dtypes(
            include="number"
        ).columns
    )

    categorical = len(
        df.select_dtypes(
            exclude="number"
        ).columns
    )

    add_metric_card(
        slide,
        Inches(2.2),
        Inches(4.2),
        Inches(3.4),
        Inches(1.4),
        "Numeric Columns",
        numeric
    )

    add_metric_card(
        slide,
        Inches(7.2),
        Inches(4.2),
        Inches(3.4),
        Inches(1.4),
        "Categorical Columns",
        categorical
    )

    add_footer(
        slide,
        3
    )


# ==========================================================
# Slide 4
# Numeric Summary
# ==========================================================

def build_numeric_slide(
    prs,
    df
):

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    add_header(
        slide,
        "Numeric KPI Summary"
    )

    add_title(
        slide,
        "Numeric Analysis"
    )

    add_subtitle(
        slide,
        "Top numerical columns with statistical summary."
    )

    table_data = build_numeric_summary(
        df
    )

    add_table_title(
        slide,
        "Statistics",
        Inches(1.8)
    )

    add_simple_table(
        slide,
        table_data,
        Inches(0.6),
        Inches(2.2),
        Inches(12.1)
    )

    add_footer(
        slide,
        4
    )
    # ==========================================================
# Bullet List Helper
# ==========================================================

def add_bullet_list(
    slide,
    title,
    items,
    left,
    top,
    width,
    height
):
    """
    Draws a title and bullet list.
    """

    title_box = slide.shapes.add_textbox(
        left,
        top,
        width,
        Inches(0.4)
    )

    tf = title_box.text_frame

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    body = slide.shapes.add_textbox(
        left,
        top + Inches(0.45),
        width,
        height
    )

    body_tf = body.text_frame
    body_tf.clear()

    if not items:
        items = [
            "No insights available."
        ]

    for index, item in enumerate(items):

        if isinstance(item, dict):

            item = (
                item.get("message")
                or item.get("text")
                or item.get("insight")
                or str(item)
            )

        if index == 0:
            paragraph = body_tf.paragraphs[0]
        else:
            paragraph = body_tf.add_paragraph()

        paragraph.text = f"• {item}"
        paragraph.level = 0
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = DARK
        paragraph.space_after = Pt(10)


# ==========================================================
# Executive Observation Helper
# ==========================================================

def build_dataset_observations(df):

    metrics = dataset_metrics(df)

    observations = []

    observations.append(
        f"Dataset contains {metrics['rows']:,} rows and {metrics['columns']} columns."
    )

    observations.append(
        f"Overall data quality score is {metrics['quality']}%."
    )

    if metrics["missing"] == 0:
        observations.append(
            "No missing values were detected."
        )
    else:
        observations.append(
            f"{metrics['missing']:,} missing values require attention."
        )

    if metrics["duplicates"] == 0:
        observations.append(
            "No duplicate records found."
        )
    else:
        observations.append(
            f"{metrics['duplicates']:,} duplicate records detected."
        )

    numeric_columns = len(
        df.select_dtypes(include="number").columns
    )

    observations.append(
        f"{numeric_columns} numeric columns available for KPI analysis."
    )

    return observations


# ==========================================================
# Slide 5
# AI Insights
# ==========================================================

def build_ai_insights_slide(
    prs,
    df,
    insights
):

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    add_header(
        slide,
        "Artificial Intelligence Insights"
    )

    add_title(
        slide,
        "AI Generated Insights"
    )

    add_subtitle(
        slide,
        "Automatically generated observations based on the uploaded dataset."
    )

    normalized = []

    if isinstance(insights, list):

        for insight in insights:

            if isinstance(insight, dict):

                normalized.append(
                    str(
                        insight.get("message")
                        or insight.get("text")
                        or insight.get("insight")
                        or insight
                    )
                )

            else:

                normalized.append(
                    str(insight)
                )

    observations = build_dataset_observations(df)

    left_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(1.8),
        Inches(5.8),
        Inches(4.8)
    )

    left_box.fill.solid()
    left_box.fill.fore_color.rgb = WHITE
    left_box.line.color.rgb = PRIMARY

    right_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.8),
        Inches(1.8),
        Inches(5.8),
        Inches(4.8)
    )

    right_box.fill.solid()
    right_box.fill.fore_color.rgb = WHITE
    right_box.line.color.rgb = PRIMARY

    add_bullet_list(
        slide,
        "AI Insights",
        normalized,
        Inches(0.8),
        Inches(2.0),
        Inches(5.2),
        Inches(4)
    )

    add_bullet_list(
        slide,
        "Dataset Observations",
        observations,
        Inches(7.0),
        Inches(2.0),
        Inches(5.0),
        Inches(4)
    )

    add_footer(
        slide,
        5
    )
    # ==========================================================
# Recommendation Helper
# ==========================================================

def normalize_recommendations(
    recommendations
):
    """
    Converts recommendation objects into plain strings.
    """

    normalized = []

    if not recommendations:
        return [
            "Monitor business KPIs regularly.",
            "Improve overall data quality.",
            "Use predictive analytics for planning."
        ]

    for item in recommendations:

        if isinstance(item, dict):

            text = (
                item.get("message")
                or item.get("text")
                or item.get("recommendation")
                or item.get("title")
                or str(item)
            )

        else:
            text = str(item)

        normalized.append(text)

    return normalized


# ==========================================================
# Slide 6
# Business Recommendations
# ==========================================================

def build_recommendation_slide(
    prs,
    recommendations
):

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    add_header(
        slide,
        "Business Recommendations"
    )

    add_title(
        slide,
        "AI Recommended Actions"
    )

    add_subtitle(
        slide,
        "Strategic recommendations generated from the uploaded dataset."
    )

    recommendation_items = normalize_recommendations(
        recommendations
    )

    recommendation_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(1.8),
        Inches(12.1),
        Inches(4.8)
    )

    recommendation_box.fill.solid()
    recommendation_box.fill.fore_color.rgb = WHITE

    recommendation_box.line.color.rgb = PRIMARY

    add_bullet_list(
        slide,
        "Recommended Business Actions",
        recommendation_items,
        Inches(0.9),
        Inches(2.0),
        Inches(11.5),
        Inches(4.2)
    )

    add_footer(
        slide,
        6
    )


# ==========================================================
# Executive Summary Helper
# ==========================================================

def build_summary_points(df):

    metrics = dataset_metrics(df)

    summary = []

    summary.append(
        f"Rows processed : {metrics['rows']:,}"
    )

    summary.append(
        f"Columns analysed : {metrics['columns']}"
    )

    summary.append(
        f"Data Quality Score : {metrics['quality']}%"
    )

    if metrics["quality"] >= 95:
        summary.append(
            "Dataset quality is Excellent."
        )

    elif metrics["quality"] >= 85:
        summary.append(
            "Dataset quality is Good."
        )

    elif metrics["quality"] >= 70:
        summary.append(
            "Dataset quality is Fair."
        )

    else:
        summary.append(
            "Dataset quality requires improvement."
        )

    numeric = len(
        df.select_dtypes(
            include="number"
        ).columns
    )

    summary.append(
        f"{numeric} numeric columns available for business analytics."
    )

    return summary


# ==========================================================
# Slide 7
# Executive Summary
# ==========================================================

def build_summary_slide(
    prs,
    df
):

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    add_header(
        slide,
        "Executive Summary"
    )

    add_title(
        slide,
        "Overall Executive Summary"
    )

    add_subtitle(
        slide,
        "A concise summary of the uploaded dataset."
    )

    metrics = dataset_metrics(df)

    score_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.8),
        Inches(3.0),
        Inches(4.8)
    )

    score_card.fill.solid()
    score_card.fill.fore_color.rgb = PRIMARY
    score_card.line.color.rgb = PRIMARY

    tf = score_card.text_frame

    p = tf.paragraphs[0]

    p.text = "Quality Score"

    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()

    p.text = f"{metrics['quality']}%"

    p.font.bold = True
    p.font.size = Pt(38)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()

    p.text = (
        "Dataset Health"
    )

    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    summary = build_summary_points(df)

    summary_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(4.1),
        Inches(1.8),
        Inches(8.5),
        Inches(4.8)
    )

    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = WHITE

    summary_box.line.color.rgb = PRIMARY

    add_bullet_list(
        slide,
        "Executive Highlights",
        summary,
        Inches(4.3),
        Inches(2.0),
        Inches(7.9),
        Inches(4.2)
    )

    add_footer(
        slide,
        7
    )
    # ==========================================================
# Slide 8
# Thank You / Closing Slide
# ==========================================================

def build_thank_you_slide(
    prs,
    filename
):

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    add_background(slide)

    # Top Banner
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.33),
        Inches(0.70)
    )

    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY
    banner.line.color.rgb = PRIMARY

    title_box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.2),
        Inches(11.5),
        Inches(0.8)
    )

    tf = title_box.text_frame

    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = (
        "Executive Business Intelligence Report"
    )
    p.font.size = Pt(22)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    info = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(2.6),
        Inches(11),
        Inches(2.2)
    )

    tf = info.text_frame

    p = tf.paragraphs[0]
    p.text = (
        "This presentation was automatically "
        "generated by the AI Business "
        "Intelligence Platform."
    )
    p.font.size = Pt(18)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = (
        f"Dataset : {filename}"
    )
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = (
        datetime.now().strftime(
            "%d %B %Y  |  %I:%M %p"
        )
    )
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Bottom Accent Bar
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        Inches(6.9),
        Inches(13.33),
        Inches(0.6)
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.color.rgb = PRIMARY

    footer = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(7.0),
        Inches(12.2),
        Inches(0.25)
    )

    tf = footer.text_frame

    p = tf.paragraphs[0]
    p.text = (
        "AI Business Intelligence Platform"
    )
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


# ==========================================================
# Helper
# ==========================================================

def ensure_output_directory(
    output_path
):

    directory = os.path.dirname(
        output_path
    )

    if directory:

        os.makedirs(
            directory,
            exist_ok=True
        )


def normalize_ai_items(items):

    if not items:
        return []

    normalized = []

    for item in items:

        if isinstance(item, dict):

            normalized.append(
                str(
                    item.get("message")
                    or item.get("text")
                    or item.get("insight")
                    or item.get("recommendation")
                    or item
                )
            )

        else:

            normalized.append(
                str(item)
            )

    return normalized
# ==========================================================
# Main Professional PowerPoint Generator
# ==========================================================

def generate_professional_ppt(
    df: pd.DataFrame,
    filename: str,
    stats: dict[str, Any] | None = None,
    insights: list[Any] | None = None,
    recommendations: list[Any] | None = None,
    output_path: str | None = None
) -> str:
    """
    Generate a complete professional PowerPoint report.

    Parameters
    ----------
    df:
        Uploaded dataset as a pandas DataFrame.

    filename:
        Original uploaded dataset filename.

    stats:
        Optional dataset statistics. Kept for compatibility
        with the Flask export route.

    insights:
        AI-generated insights. Values may be strings or
        dictionaries.

    recommendations:
        AI-generated business recommendations. Values may
        be strings or dictionaries.

    output_path:
        Full location where the .pptx file will be saved.

    Returns
    -------
    str
        The saved PowerPoint file path.
    """

    if not isinstance(
        df,
        pd.DataFrame
    ):
        raise TypeError(
            "df must be a pandas DataFrame."
        )

    if df.empty:
        raise ValueError(
            "Cannot generate a PowerPoint report "
            "from an empty dataset."
        )

    if not filename:
        filename = "uploaded_dataset"

    if not output_path:
        safe_name = os.path.splitext(
            os.path.basename(
                filename
            )
        )[0]

        output_path = os.path.join(
            "exports",
            (
                f"{safe_name}_"
                "professional_report.pptx"
            )
        )

    if not output_path.lower().endswith(
        ".pptx"
    ):
        output_path = (
            f"{output_path}.pptx"
        )

    ensure_output_directory(
        output_path
    )

    normalized_insights = (
        normalize_ai_items(
            insights
        )
    )

    normalized_recommendations = (
        normalize_ai_items(
            recommendations
        )
    )

    if not normalized_insights:
        normalized_insights = (
            build_dataset_observations(
                df
            )
        )

    if not normalized_recommendations:
        normalized_recommendations = [
            (
                "Monitor the most important "
                "business KPIs regularly."
            ),
            (
                "Clean missing and duplicate "
                "records before critical analysis."
            ),
            (
                "Use forecasting models for "
                "forward-looking business planning."
            ),
            (
                "Validate unusual values and "
                "outliers before making decisions."
            )
        ]

    try:
        prs = create_presentation()

        # Slide 1
        build_cover_slide(
            prs=prs,
            filename=filename
        )

        # Slide 2
        build_overview_slide(
            prs=prs,
            df=df
        )

        # Slide 3
        build_kpi_slide(
            prs=prs,
            df=df
        )

        # Slide 4
        build_numeric_slide(
            prs=prs,
            df=df
        )

        # Slide 5
        build_ai_insights_slide(
            prs=prs,
            df=df,
            insights=normalized_insights
        )

        # Slide 6
        build_recommendation_slide(
            prs=prs,
            recommendations=(
                normalized_recommendations
            )
        )

        # Slide 7
        build_summary_slide(
            prs=prs,
            df=df
        )

        # Slide 8
        build_thank_you_slide(
            prs=prs,
            filename=filename
        )

        prs.save(
            output_path
        )

    except PermissionError as error:
        raise PermissionError(
            (
                "PowerPoint file could not be saved. "
                "Close the existing presentation if "
                "it is currently open and try again."
            )
        ) from error

    except Exception as error:
        raise RuntimeError(
            (
                "Professional PowerPoint generation "
                f"failed: {str(error)}"
            )
        ) from error

    if not os.path.exists(
        output_path
    ):
        raise FileNotFoundError(
            (
                "PowerPoint generation completed "
                "without creating the output file."
            )
        )

    return output_path