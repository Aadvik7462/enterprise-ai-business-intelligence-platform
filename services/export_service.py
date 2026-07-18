import os
from datetime import datetime

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle
)


def generate_pdf_report(filename, stats, insights, recommendations, export_folder):
    os.makedirs(export_folder, exist_ok=True)

    report_name = f"Executive_Report_{filename.rsplit('.', 1)[0]}.pdf"
    report_path = os.path.join(export_folder, report_name)

    doc = SimpleDocTemplate(
        report_path,
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("AI Business Intelligence Platform", styles["Title"]))
    story.append(Paragraph("Executive PDF Report", styles["Heading2"]))
    story.append(Spacer(1, 12))

    story.append(Paragraph(f"<b>Dataset:</b> {filename}", styles["Normal"]))
    story.append(Paragraph(f"<b>Generated:</b> {datetime.now().strftime('%d-%m-%Y %I:%M %p')}", styles["Normal"]))
    story.append(Spacer(1, 18))

    story.append(Paragraph("Dataset Overview", styles["Heading2"]))

    stats_data = [
        ["Metric", "Value"],
        ["Total Rows", stats["rows"]],
        ["Total Columns", stats["columns"]],
        ["Missing Values", stats["missing_values"]],
        ["Duplicate Rows", stats["duplicate_rows"]],
        ["Numeric Columns", stats.get("numeric_columns", "-")],
        ["Categorical Columns", stats.get("categorical_columns", "-")],
        ["Memory Usage", f"{stats['memory_usage']} MB"],
        ["Quality Score", f"{stats['quality_score']}%"]
    ]

    stats_table = Table(stats_data, colWidths=[220, 220])
    stats_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0f172a")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#f8fafc")),
        ("PADDING", (0, 0), (-1, -1), 9),
    ]))

    story.append(stats_table)
    story.append(Spacer(1, 18))

    story.append(Paragraph("AI Generated Insights", styles["Heading2"]))

    for insight in insights[:8]:
        story.append(Paragraph(f"- {insight}", styles["Normal"]))
        story.append(Spacer(1, 6))

    story.append(Spacer(1, 12))
    story.append(Paragraph("AI Recommendations", styles["Heading2"]))

    for rec in recommendations:
        story.append(Paragraph(f"<b>{rec['title']}:</b> {rec['message']}", styles["Normal"]))
        story.append(Spacer(1, 6))

    story.append(Spacer(1, 18))
    story.append(Paragraph("Report Summary", styles["Heading2"]))

    summary_text = (
        "This report provides an automated overview of the uploaded dataset, "
        "including data quality, missing values, duplicate records, AI insights, "
        "and recommendations for further analysis."
    )

    story.append(Paragraph(summary_text, styles["Normal"]))

    doc.build(story)

    return report_name
from pptx import Presentation
from pptx.util import Inches, Pt


def generate_ppt_report(filename, stats, insights, recommendations, export_folder):
    os.makedirs(export_folder, exist_ok=True)

    ppt_name = f"Executive_Report_{filename.rsplit('.', 1)[0]}.pptx"
    ppt_path = os.path.join(export_folder, ppt_name)

    prs = Presentation()

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "AI Business Intelligence Platform"
        slide.placeholders[1].text = f"Executive Report\nDataset: {filename}"

    def add_bullet_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()

        for item in bullets:
            p = body.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(18)

    add_title_slide()

    overview = [
        f"Total Rows: {stats['rows']}",
        f"Total Columns: {stats['columns']}",
        f"Missing Values: {stats['missing_values']}",
        f"Duplicate Rows: {stats['duplicate_rows']}",
        f"Quality Score: {stats['quality_score']}%",
        f"Memory Usage: {stats['memory_usage']} MB",
    ]
    add_bullet_slide("Dataset Overview", overview)

    add_bullet_slide("AI Generated Insights", insights[:6])

    rec_bullets = [
        f"{rec['title']}: {rec['message']}"
        for rec in recommendations
    ]
    add_bullet_slide("AI Recommendations", rec_bullets)

    add_bullet_slide("Conclusion", [
        "Dataset analysis completed successfully.",
        "Use the Executive Dashboard for interactive exploration.",
        "Use PDF/PPT exports for business reporting.",
    ])

    prs.save(ppt_path)

    return ppt_name
import shutil
import zipfile


def generate_export_package(
    filename,
    source_filepath,
    stats,
    insights,
    recommendations,
    export_folder
):
    """
    Generate PDF and PowerPoint reports and package them with
    the uploaded dataset inside a ZIP file.
    """

    os.makedirs(export_folder, exist_ok=True)

    dataset_name = os.path.splitext(filename)[0]

    package_folder = os.path.join(
        export_folder,
        f"Export_Package_{dataset_name}"
    )

    # Remove an older temporary package folder
    if os.path.exists(package_folder):
        shutil.rmtree(package_folder)

    os.makedirs(package_folder, exist_ok=True)

    # Generate PDF inside the temporary package folder
    pdf_name = generate_pdf_report(
        filename=filename,
        stats=stats,
        insights=insights,
        recommendations=recommendations,
        export_folder=package_folder
    )

    # Generate PowerPoint inside the temporary package folder
    ppt_name = generate_ppt_report(
        filename=filename,
        stats=stats,
        insights=insights,
        recommendations=recommendations,
        export_folder=package_folder
    )

    # Copy uploaded or cleaned dataset into the package
    copied_dataset_path = os.path.join(
        package_folder,
        os.path.basename(filename)
    )

    shutil.copy2(source_filepath, copied_dataset_path)

    zip_name = f"AI_BI_Export_Package_{dataset_name}.zip"
    zip_path = os.path.join(export_folder, zip_name)

    # Remove the previous ZIP if it exists
    if os.path.exists(zip_path):
        os.remove(zip_path)

    # Build ZIP
    with zipfile.ZipFile(
        zip_path,
        mode="w",
        compression=zipfile.ZIP_DEFLATED
    ) as zip_file:

        files_to_package = [
            pdf_name,
            ppt_name,
            os.path.basename(filename)
        ]

        for packaged_file in files_to_package:
            full_path = os.path.join(package_folder, packaged_file)

            if os.path.exists(full_path):
                zip_file.write(
                    full_path,
                    arcname=packaged_file
                )

    # Delete temporary folder after ZIP creation
    shutil.rmtree(package_folder)

    return zip_name